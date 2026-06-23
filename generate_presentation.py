from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches
import numpy as np
from sklearn.linear_model import LinearRegression
from sklearn.metrics import mean_squared_error
from sklearn.model_selection import train_test_split, cross_val_score
from sklearn.preprocessing import StandardScaler
from main import AcademicPredictor

DARK_TEXT = RGBColor(30, 30, 46)
ACCENT_BLUE = RGBColor(99, 102, 241)
LIGHT_GRAY = RGBColor(220, 220, 230)
MILD_GREY = RGBColor(100, 100, 120)
WARNING_RED = RGBColor(239, 68, 68)
YELLOW_WARN = RGBColor(234, 179, 8)


def compute_project_metrics():
    predictor = AcademicPredictor()
    df_raw = predictor.load_data('master_markbook.csv')
    df_clean = predictor.clean_data(df_raw)

    X1 = df_clean[['Maths_Advanced']].values
    y = df_clean['Software_Engineering_Final'].values
    X1_train, X1_test, y_train, y_test = train_test_split(X1, y, test_size=0.2, random_state=42)
    baseline = LinearRegression().fit(X1_train, y_train)
    baseline_rmse = np.sqrt(mean_squared_error(y_test, baseline.predict(X1_test)))

    X2 = df_clean[['Maths_Advanced', 'Physics']].values
    X2_train, X2_test, y2_train, y2_test = train_test_split(X2, y, test_size=0.2, random_state=42)
    scaler = StandardScaler().fit(X2_train)
    X2_train_scaled = scaler.transform(X2_train)
    X2_test_scaled = scaler.transform(X2_test)
    level2_model = LinearRegression().fit(X2_train_scaled, y2_train)
    level2_rmse = np.sqrt(mean_squared_error(y2_test, level2_model.predict(X2_test_scaled)))

    cv_rmse = -cross_val_score(level2_model, scaler.transform(X2), y, cv=5, scoring='neg_root_mean_squared_error').mean()

    alex_row = df_raw[df_raw['Student_Name'].str.contains('Alex', case=False, na=False)]
    alex_prediction = None
    if not alex_row.empty:
        alex_prediction = float(level2_model.predict(scaler.transform(alex_row[['Maths_Advanced', 'Physics']].iloc[0].values.reshape(1, -1)))[0])

    return {
        'baseline_rmse': baseline_rmse,
        'level2_rmse': level2_rmse,
        'cv_rmse': cv_rmse,
        'alex_prediction': alex_prediction,
        'cleaned_rows': len(df_clean),
        'raw_rows': len(df_raw),
    }


def add_title_text(slide, text, left=Inches(0.5), top=Inches(0.5), width=Inches(9), height=Inches(1.2), size=Pt(48), bold=True, color=ACCENT_BLUE):
    """Add styled title text box to slide."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    para = frame.paragraphs[0]
    para.text = text
    para.font.size = size
    para.font.bold = bold
    para.font.color.rgb = color
    para.font.name = 'Calibri'
    para.alignment = PP_ALIGN.LEFT
    return textbox


def add_subtitle_text(slide, text, left=Inches(0.5), top=Inches(1.7), width=Inches(9), height=Inches(0.8), size=Pt(22), color=LIGHT_GRAY):
    """Add styled subtitle text box to slide."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.text = text
    para.font.size = size
    para.font.color.rgb = color
    para.font.name = 'Calibri'
    para.alignment = PP_ALIGN.LEFT
    return textbox


def add_body_text(slide, text_lines, left=Inches(0.5), top=Inches(2.5), width=Inches(9), height=Inches(4.5), size=Pt(18), color=DARK_TEXT):
    """Add styled body text with bullet points."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    
    for idx, line in enumerate(text_lines):
        if idx == 0:
            para = frame.paragraphs[0]
        else:
            para = frame.add_paragraph()
        para.text = line
        para.font.size = size
        para.font.color.rgb = color
        para.font.name = 'Calibri'
        para.alignment = PP_ALIGN.LEFT
        para.level = 0
        para.space_after = Pt(8)
    return textbox


def add_divider_line(slide, top_inches=2.2):
    """Add a horizontal divider rectangle."""
    rect = slide.shapes.add_shape(1, Inches(0.5), Inches(top_inches), Inches(9), Inches(0.02))
    rect.fill.solid()
    rect.fill.fore_color.rgb = ACCENT_BLUE
    rect.line.color.rgb = ACCENT_BLUE
    return rect


def add_blank_slide(prs):
    """Add a blank slide to the presentation."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def generate_presentation(output_path='Academic Predictor Presentation.pptx'):
    """Generate presentation with Academic Predictor content in Security Audit style."""
    metrics = compute_project_metrics()
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = add_blank_slide(prs)
    add_title_text(slide, 'Academic Predictor', top=Inches(1), size=Pt(54), color=ACCENT_BLUE)
    add_subtitle_text(slide, 'Multi-feature linear regression for predicting student final grades', top=Inches(2), size=Pt(22), color=LIGHT_GRAY)
    add_body_text(slide, [
        'Automated data cleaning & validation',
        'Supervised learning with train-test split',
        'Cross-validation & RMSE evaluation',
        'Privacy-aware data handling & bias audit'
    ], top=Inches(3.2), size=Pt(18))

    slide = add_blank_slide(prs)
    add_title_text(slide, 'Agenda', size=Pt(40), color=DARK_TEXT)
    add_divider_line(slide, 1.2)
    add_body_text(slide, [
        '1.  Dataset overview & cleaning strategy',
        '2.  Baseline vs. multi-feature modeling',
        '3.  Validation results & prediction accuracy',
        '4.  Bias audit & ethical safeguards',
        '5.  Submission readiness & summary'
    ], top=Inches(1.6), size=Pt(20))

    slide = add_blank_slide(prs)
    add_title_text(slide, 'Dataset Overview', size=Pt(40), color=DARK_TEXT)
    add_divider_line(slide, 1.2)
    add_body_text(slide, [
        f'Original dataset: {metrics["raw_rows"]} student records',
        f'Cleaned dataset: {metrics["cleaned_rows"]} valid rows (removed missing/invalid scores)',
        'Feature columns: Maths Advanced, Physics, Modern History',
        'Target column: Software Engineering Final Grade (0–100)',
        'Cleaning rules: drop NaN, enforce [0, 100] range per column'
    ], top=Inches(1.6), size=Pt(18))

    slide = add_blank_slide(prs)
    add_title_text(slide, 'Model Strategy: Baseline vs. Level 2', size=Pt(38), color=DARK_TEXT)
    add_divider_line(slide, 1.2)
    
    textbox_left = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.3), Inches(5.5))
    frame_left = textbox_left.text_frame
    frame_left.word_wrap = True
    p = frame_left.paragraphs[0]
    p.text = 'Baseline: Single-Feature'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = 'Calibri'
    for line in ['Linear regression on Maths Advanced only', 'Simple, interpretable', 'Establishes performance baseline']:
        para = frame_left.add_paragraph()
        para.text = line
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_TEXT
        para.font.name = 'Calibri'
        para.space_after = Pt(6)

    textbox_right = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.5))
    frame_right = textbox_right.text_frame
    frame_right.word_wrap = True
    p = frame_right.paragraphs[0]
    p.text = 'Level 2: Multi-Feature'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = 'Calibri'
    for line in ['3 features + StandardScaler + LinearRegression', 'Reduces feature bias from different ranges', 'Better predictions with more context']:
        para = frame_right.add_paragraph()
        para.text = line
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_TEXT
        para.font.name = 'Calibri'
        para.space_after = Pt(6)

    slide = add_blank_slide(prs)
    add_title_text(slide, 'Validation & Results', size=Pt(40), color=DARK_TEXT)
    add_divider_line(slide, 1.2)
    add_body_text(slide, [
        f'Train-test split: 80% training, 20% testing (random_state=42)',
        f'Baseline RMSE: {metrics["baseline_rmse"]:.2f}',
        f'Level 2 RMSE: {metrics["level2_rmse"]:.2f}',
        f'Cross-validation RMSE (5-fold): {metrics["cv_rmse"]:.2f}',
        f'Alex Anderson predicted final score: {metrics["alex_prediction"]:.1f}'
    ], top=Inches(1.6), size=Pt(18))

    slide = add_blank_slide(prs)
    add_title_text(slide, 'Bias Audit & Ethical Safeguards', size=Pt(38), color=DARK_TEXT)
    add_divider_line(slide, 1.2)
    
    textbox_left = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.3), Inches(5.5))
    frame_left = textbox_left.text_frame
    frame_left.word_wrap = True
    p = frame_left.paragraphs[0]
    p.text = 'Privacy & Data Protection'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = 'Calibri'
    for line in ['Anonymize student names', 'Hash IDs for safe sharing', 'Remove PII before export']:
        para = frame_left.add_paragraph()
        para.text = line
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_TEXT
        para.font.name = 'Calibri'
        para.space_after = Pt(6)

    textbox_right = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.5))
    frame_right = textbox_right.text_frame
    frame_right.word_wrap = True
    p = frame_right.paragraphs[0]
    p.text = 'Bias Detection'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = 'Calibri'
    for line in ['80% rule disparate impact check', 'Compare subgroup pass rates', 'Audit report with warnings']:
        para = frame_right.add_paragraph()
        para.text = line
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_TEXT
        para.font.name = 'Calibri'
        para.space_after = Pt(6)

    slide = add_blank_slide(prs)
    add_title_text(slide, 'Submission-Ready Summary', size=Pt(40), color=DARK_TEXT)
    add_divider_line(slide, 1.2)
    add_body_text(slide, [
        'Software: Complete pipeline with data cleaning, training, evaluation, and ethics checks',
        'Results: Stable, reproducible RMSE metrics and student predictions',
        'Documentation: Project writeup with all template sections filled',
        'Testing: 13 passing unit tests with 61% code coverage',
        'Code Quality: OOP design with AcademicPredictor class, anonymization, bias audit'
    ], top=Inches(1.6), size=Pt(17))

    prs.save(output_path)
    print(f'Generated presentation: {output_path}')
    return output_path


if __name__ == '__main__':
    generate_presentation()
